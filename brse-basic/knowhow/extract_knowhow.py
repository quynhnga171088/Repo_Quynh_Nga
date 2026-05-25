# -*- coding: utf-8 -*-
"""Quét workspace và append text vào các file Knowhow-*.md"""
from __future__ import annotations

import html
import os
import re
import sys
import traceback
import zipfile
import xml.etree.ElementTree as ET
from html.parser import HTMLParser
from pathlib import Path

ROOT = Path(r"c:\Dự án\Toàn bộ file tài nguyên của Nga")
OUT_DIR = ROOT / "knowhow"
CHAR_LIMIT = 800_000

BUCKETS = {
    "phan_tich_nghiep_vu": "Knowhow-phan-tich-nghiep-vu",
    "quan_ly": "Knowhow-quan-ly",
    "tieng_nhat": "Knowhow-tieng-nhat",
    "ai": "Knowhow-ai",
    "content": "Knowhow-content",
    "viet_sach": "Knowhow-viet-sach",
}

# Prefix path (Windows \\) -> buckets mặc định; khớp prefix dài nhất trước
FOLDER_RULES: list[tuple[str, list[str]]] = [
    ("Udemy\\15-8-2025", ["ai", "tieng_nhat", "content"]),
    ("Udemy\\Udemy - N1 Course", ["tieng_nhat"]),
    ("Udemy\\Udemy - 8-29", ["tieng_nhat", "content"]),
    ("Udemy\\Udemy - 7-31", ["tieng_nhat", "phan_tich_nghiep_vu", "ai"]),
    ("Udemy\\Udemy - 4-30", ["tieng_nhat", "content"]),
    ("Udemy", ["content", "tieng_nhat"]),
    ("Đào tạo IT Comtor", ["tieng_nhat", "quan_ly"]),
    ("Kinh nghiệm đào tạo\\Đào tạo BrSE noncode", ["phan_tich_nghiep_vu", "ai", "quan_ly"]),
    ("Kinh nghiệm đào tạo\\Đào tạo IT Comtor", ["tieng_nhat", "quan_ly"]),
    ("Kinh nghiệm đào tạo", ["quan_ly", "content", "ai"]),
    ("Kinh nghiệm phỏng vấn ứng viên", ["tieng_nhat", "phan_tich_nghiep_vu", "quan_ly"]),
    ("Kinh nghiệm PM", ["quan_ly", "tieng_nhat"]),
    ("Kinh nghiệm test", ["phan_tich_nghiep_vu"]),
    ("Kỹ năng", ["content", "quan_ly"]),
    ("Trợ lý dự án BrSE\\Xây dựng Website LMS", ["phan_tich_nghiep_vu", "ai"]),
    ("Trợ lý dự án BrSE\\Project 1-BrSE Foundation\\1. Bài viết", [
        "phan_tich_nghiep_vu", "content", "tieng_nhat", "ai", "viet_sach",
    ]),
    ("Trợ lý dự án BrSE\\Project 1-BrSE Foundation\\6. Workshop", ["content", "viet_sach"]),
    ("Trợ lý dự án BrSE\\Project 1-BrSE Foundation\\4. Tài liệu xây dựng khoá học", [
        "phan_tich_nghiep_vu", "content", "ai",
    ]),
    ("Trợ lý dự án BrSE\\Project 1-BrSE Foundation\\2. Slide", ["phan_tich_nghiep_vu", "ai"]),
    ("Trợ lý dự án BrSE", ["phan_tich_nghiep_vu", "content", "quan_ly"]),
]
FOLDER_RULES.sort(key=lambda x: len(x[0]), reverse=True)

KEYWORD_BUCKETS: dict[str, list[str]] = {
    "tieng_nhat": [
        "tiếng nhật", "tieng nhat", "comtor", "n1 course", " n1", " n2", "dịch",
        "mail tiếng", "mail chuẩn", "日本", "nhật bản", "議事録", "秘密保持", "翻訳",
    ],
    "ai": [
        " ai", "chatgpt", "cursor", "copilot", "prompt", "notebooklm", "claude",
        "dify", "gpt", " ocr", "ứng dụng ai",
    ],
    "content": [
        "content", "bài viết", "workshop", "truyền thông", "portfolio", "udemy",
        "đăng fb", "kênh", "ý tưởng bài",
    ],
    "viet_sach": ["sách", "workbook", "insight-report", "nháp", "xuất bản"],
    "quan_ly": [
        "pm", "quản lý", "project management", "weekly report", "lesson learn",
        "meeting minutes", "jira", "task-format",
    ],
    "phan_tich_nghiep_vu": [
        "brse", "nghiệp vụ", "requirement", "phân tích", "điều tra",
        "khung phân tích", "prd", "spec", "testcase", "erd",
    ],
}

SKIP_EXT = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".mp4", ".mov", ".avi", ".mkv"}
TEXT_EXT = {
    ".md", ".txt", ".html", ".htm", ".docx", ".doc", ".xlsx", ".xlsm", ".xls",
    ".pptx", ".pdf", ".drawio", ".csv",
}
PDF_MIN_CHARS = 80


class _HTMLText(HTMLParser):
    def __init__(self):
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        t = data.strip()
        if t:
            self.parts.append(t)

    def get_text(self) -> str:
        return "\n".join(self.parts)


def norm_path(p: str) -> str:
    return p.replace("/", "\\")


def classify(rel_path: str) -> set[str]:
    buckets: set[str] = set()
    for prefix, bs in FOLDER_RULES:
        if rel_path.startswith(prefix):
            buckets.update(bs)
            break
    hay = rel_path.lower()
    for b, kws in KEYWORD_BUCKETS.items():
        for kw in kws:
            if kw in hay:
                buckets.add(b)
                break
    if not buckets:
        buckets.add("content")
    return buckets


class Writers:
    def __init__(self):
        self._paths: dict[str, Path] = {}
        self._sizes: dict[str, int] = {}
        self._part: dict[str, int] = {}

    def _base_path(self, bucket: str) -> Path:
        if bucket not in self._paths:
            name = BUCKETS[bucket]
            self._paths[bucket] = OUT_DIR / f"{name}.md"
            self._sizes[bucket] = 0
            self._part[bucket] = 1
        return self._paths[bucket]

    def append(self, bucket: str, block: str) -> None:
        path = self._base_path(bucket)
        size = self._sizes[bucket]
        if size + len(block) > CHAR_LIMIT and size > 0:
            self._part[bucket] += 1
            n = self._part[bucket]
            name = BUCKETS[bucket]
            path = OUT_DIR / f"{name}-part-{n:02d}.md"
            self._paths[bucket] = path
            self._sizes[bucket] = 0
            header = f"# {name} (part {n:02d})\n\n"
            path.write_text(header, encoding="utf-8")
            self._sizes[bucket] = len(header)
        with path.open("a", encoding="utf-8") as f:
            f.write(block)
        self._sizes[bucket] += len(block)


writers = Writers()
skipped: list[tuple[str, str]] = []
stats = {"processed": 0, "empty": 0, "errors": 0}


def log_skip(rel: str, reason: str) -> None:
    skipped.append((rel, reason))


def extract_html(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="replace")
    p = _HTMLText()
    p.feed(raw)
    return p.get_text()


def extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        parts.append(f"## Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts)


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            parts.append(f"## Slide {i}")
            parts.extend(slide_text)
    return "\n".join(parts)


def extract_pdf(path: Path) -> tuple[str | None, str | None]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts: list[str] = []
    for page in reader.pages:
        t = page.extract_text() or ""
        if t.strip():
            parts.append(t)
    text = "\n".join(parts).strip()
    if len(text) < PDF_MIN_CHARS:
        return None, "pdf-scan-or-empty"
    return text, None


def extract_drawio(path: Path) -> str:
    with zipfile.ZipFile(path, "r") as zf:
        xml = zf.read("diagram.xml" if "diagram.xml" in zf.namelist() else zf.namelist()[0])
    root = ET.fromstring(xml)
    texts: list[str] = []
    for el in root.iter():
        if el.text and el.text.strip():
            texts.append(el.text.strip())
        val = el.attrib.get("value")
        if val and val.strip():
            texts.append(html.unescape(val.strip()))
    return "\n".join(texts)


def extract_file(path: Path, rel: str) -> tuple[str | None, str | None]:
    ext = path.suffix.lower()
    try:
        if ext in {".md", ".txt", ".csv"}:
            return path.read_text(encoding="utf-8", errors="replace"), None
        if ext in {".html", ".htm"}:
            return extract_html(path), None
        if ext == ".docx":
            return extract_docx(path), None
        if ext in {".xlsx", ".xlsm"}:
            return extract_xlsx(path), None
        if ext == ".pptx":
            return extract_pptx(path), None
        if ext == ".pdf":
            return extract_pdf(path)
        if ext == ".drawio":
            return extract_drawio(path), None
        if ext == ".doc":
            return None, "doc-legacy-unsupported"
        return None, f"unsupported-ext-{ext}"
    except Exception as e:
        err = str(e).lower()
        if "docx" in err or "zip" in err or "font" in err:
            return None, "docx-error"
        return None, f"extract-error: {e}"


def format_block(rel: str, text: str) -> str:
    return f"\n\n---\n**Nguồn:** `{rel}`\n\n{text}\n"


def process_file(path: Path) -> None:
    rel = norm_path(os.path.relpath(path, ROOT))
    ext = path.suffix.lower()

    if "knowhow" in rel.split("\\"):
        return

    if ext in SKIP_EXT:
        log_skip(rel, "image-or-video")
        return

    if ext not in TEXT_EXT:
        log_skip(rel, f"unsupported-ext{ext}")
        return

    text, err = extract_file(path, rel)
    if err:
        log_skip(rel, err)
        return
    if not text or not text.strip():
        stats["empty"] += 1
        log_skip(rel, "empty-content")
        return

    buckets = classify(rel)
    block = format_block(rel, text.strip())
    for b in buckets:
        writers.append(b, block)
    stats["processed"] += 1


def write_skipped() -> None:
    lines = [
        "# Knowhow — File bỏ qua / không extract được",
        "",
        "> Cập nhật tự động bởi `extract_knowhow.py`.",
        "",
        "| Đường dẫn | Lý do |",
        "|-----------|--------|",
    ]
    for rel, reason in sorted(skipped):
        lines.append(f"| `{rel}` | {reason} |")
    (OUT_DIR / "Knowhow-skipped-files.md").write_text("\n".join(lines) + "\n", encoding="utf-8")


def write_map() -> None:
    lines = [
        "# Map folder → bộ knowhow",
        "",
        "Quy tắc: prefix dài nhất khớp trước + keyword trong đường dẫn/tên file.",
        "",
        "## Folder rules",
        "",
    ]
    for prefix, bs in sorted(FOLDER_RULES, key=lambda x: len(x[0]), reverse=True):
        lines.append(f"- `{prefix}` → {', '.join(bs)}")
    lines.extend(["", "## Keyword buckets", ""])
    for b, kws in KEYWORD_BUCKETS.items():
        lines.append(f"- **{b}**: {', '.join(kws[:8])}…")
    (OUT_DIR / "Knowhow-folder-map.md").write_text("\n".join(lines) + "\n", encoding="utf-8")


def main() -> int:
    # Xóa nội dung cũ (giữ header) — chỉ các file part nếu có từ lần chạy trước
    for p in OUT_DIR.glob("Knowhow-*-part-*.md"):
        p.unlink()
    for bucket, name in BUCKETS.items():
        p = OUT_DIR / f"{name}.md"
        title = {
            "phan_tich_nghiep_vu": "Phân tích nghiệp vụ",
            "quan_ly": "Quản lý",
            "tieng_nhat": "Tiếng Nhật",
            "ai": "AI",
            "content": "Content",
            "viet_sach": "Viết sách",
        }[bucket]
        p.write_text(
            f"# Knowhow — {title}\n\n"
            f"> File tổng hợp text từ workspace.\n\n",
            encoding="utf-8",
        )

    dirs_sorted = []
    for dirpath, dirnames, filenames in os.walk(ROOT):
        if "knowhow" in Path(dirpath).parts:
            continue
        dirs_sorted.append((dirpath, sorted(filenames)))

    for dirpath, filenames in dirs_sorted:
        for fn in filenames:
            process_file(Path(dirpath) / fn)

    write_skipped()
    write_map()

    summary_path = OUT_DIR / "_extract-summary.json"
    import json

    sizes = {}
    for bucket, name in BUCKETS.items():
        total = 0
        for p in OUT_DIR.glob(f"{name}*.md"):
            if p.name == "Knowhow-skipped-files.md":
                continue
            total += p.stat().st_size
        sizes[bucket] = total

    summary = {"stats": stats, "skipped_count": len(skipped), "sizes_bytes": sizes}
    summary_path.write_text(json.dumps(summary, indent=2), encoding="utf-8")
    print(json.dumps(summary, indent=2))
    return 0


if __name__ == "__main__":
    sys.exit(main())
